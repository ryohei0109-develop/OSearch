"""PowerPoint ファイルのパーサー"""

from pptx import Presentation


def parse(file_path: str) -> list[dict]:
    """
    .pptx ファイルを解析してコンテンツリストを返す。
    各シェイプのテキストをスライド番号・シェイプIDとともに返す。
    """
    contents = []
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"[警告] PowerPointファイルの読み込みに失敗しました: {file_path} ({e})")
        return contents

    for slide_num, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            contents.append({
                "location": {
                    "slide": slide_num,
                    "shape_id": shape.shape_id,
                    "sheet": None,
                    "cell": None,
                    "page": None,
                },
                "content": text,
            })

    return contents
